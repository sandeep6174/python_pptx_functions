from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import MSO_AUTO_SIZE
from pptx.enum.shapes import MSO_AUTO_SHAPE_TYPE
from pptx.enum.text import PP_PARAGRAPH_ALIGNMENT


"""Shapes dictionary contains different shapes used in power point presentation. 
We are using this dictionary in add_shape function"""

shapes_dict = {'rectangle': MSO_AUTO_SHAPE_TYPE.RECTANGLE, 'oval': MSO_AUTO_SHAPE_TYPE.OVAL,
               'circle': MSO_AUTO_SHAPE_TYPE.OVAL,
               'right triangle': MSO_AUTO_SHAPE_TYPE.RIGHT_TRIANGLE,
               'rounded rectangle': MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, 'arc': MSO_AUTO_SHAPE_TYPE.ARC,
               'balloon': MSO_AUTO_SHAPE_TYPE.BALLOON,
               'cloud': MSO_AUTO_SHAPE_TYPE.CLOUD, 'can': MSO_AUTO_SHAPE_TYPE.CAN,
               'donut': MSO_AUTO_SHAPE_TYPE.DONUT, 'diamond': MSO_AUTO_SHAPE_TYPE.DIAMOND,
               'explosion':  MSO_AUTO_SHAPE_TYPE.EXPLOSION1, 'funnel': MSO_AUTO_SHAPE_TYPE.FUNNEL,
               'gear': MSO_AUTO_SHAPE_TYPE.GEAR_6, 'heart': MSO_AUTO_SHAPE_TYPE.HEART,
               'heptagon':  MSO_AUTO_SHAPE_TYPE.HEPTAGON, 'hexagon': MSO_AUTO_SHAPE_TYPE.HEXAGON,
               'moon': MSO_AUTO_SHAPE_TYPE.MOON, 'pentagon': MSO_AUTO_SHAPE_TYPE.PENTAGON,
               'pie':  MSO_AUTO_SHAPE_TYPE.PIE, 'octagon': MSO_AUTO_SHAPE_TYPE.OCTAGON,
               'smiley face': MSO_AUTO_SHAPE_TYPE.SMILEY_FACE, 'star': MSO_AUTO_SHAPE_TYPE.STAR_10_POINT,
               'sun': MSO_AUTO_SHAPE_TYPE.SUN, 'tear': MSO_AUTO_SHAPE_TYPE.TEAR,
               'trapezoid': MSO_AUTO_SHAPE_TYPE.TRAPEZOID, 'wave': MSO_AUTO_SHAPE_TYPE.WAVE,
               'up arrow': MSO_AUTO_SHAPE_TYPE.UP_ARROW, 'down arrow':  MSO_AUTO_SHAPE_TYPE.DOWN_ARROW,
               'left arrow':  MSO_AUTO_SHAPE_TYPE.LEFT_ARROW,
               'right arrow': MSO_AUTO_SHAPE_TYPE.RIGHT_ARROW, 'u turn arrow':  MSO_AUTO_SHAPE_TYPE.U_TURN_ARROW}

"""This function takes width and height parameters to determine slide size and returns a ppt object
from main Presentation function of python-pptx library.Default size is Standard 4:3"""


def new_presentation(width=12, height=9):
    ppt = Presentation()
    ppt.slide_width = Inches(width)
    ppt.slide_height = Inches(height)
    return ppt


"""This function helps to add new slide to presentation and returns slide object to further operate on the slide. 
It takes ppt object, slide format as input parameters. Default slide format is blank"""


def add_slide(ppt, slide_format=6):
    return ppt.slides.add_slide(ppt.slide_layouts[slide_format])


"""Adds new shape to slide.Takes slide object,shape name, position, color boarder size and color of shape and 
filling color of shape as input parameters. 
All values are in inches and for colors it takes rgb value in form of tuple."""


def add_shape(slide, shape, left=0, top=0, width=0, height=0, shape_color=(0, 0, 0), boarder=0,
              boarder_color=(0, 0, 0)):
    shape = shape.lower()
    top, left, width, height = Inches(top), Inches(left), Inches(width), Inches(height)
    shapes = slide.shapes.add_shape(shapes_dict[shape], left, top, width, height)
    if boarder != 0:
        line = shapes.line
        line.color.rgb = RGBColor(boarder_color[0], boarder_color[1], boarder_color[2])
        line.width = Inches(boarder)
    fill = shapes.fill
    fill.solid()
    fill.fore_color.rgb = RGBColor(shape_color[0], shape_color[1], shape_color[2])


"""Adds text to slide. Takes slide object, text, position and font properties like size,name, color, bold, italic,
underlined as parameters. If you want to add hyperlink make hyperlink parameter to true and provide hyperlink address. 
Default font color is black and hyperlink url is google.com"""


def add_text(slide, text, left=0, top=0, width=0, height=0, font_size=16, font_name='Arial', color=(0, 0, 0),
             bold=False, italic=False, underline=False, hyperlink=False, hyperlink_address='https://www.google.com/'):
    left, top, width, height = Inches(left), Inches(top), Inches(width), Inches(height)
    tf = slide.shapes.add_textbox(left, top, width, height).text_frame
    if hyperlink:
        run = tf.paragraphs[0].add_run()
        run.text = text
        run.hyperlink.address = hyperlink_address
    else:
        tf.text = text
    tf.paragraphs[0].font.bold = bold
    tf.paragraphs[0].font.italic = italic
    tf.paragraphs[0].font.underline = underline
    tf.paragraphs[0].font.color.rgb = RGBColor(color[0], color[1], color[2])
    tf.paragraphs[0].font.size = Pt(font_size)
    tf.paragraphs[0].font.name = font_name.capitalize()


"""Create a new table and returns a table object to work on cells in table"""


def add_table(slide, rows=0, columns=0, left=0, top=0, width=0, height=0):
    left, top, width, height = Inches(left), Inches(top), Inches(width), Inches(height)
    return slide.shapes.add_table(rows, columns, left, top, width, height).table


"""return all cells in table"""


def iter_cells(table):
    for row in table.rows:
        for cell in row.cells:
            yield cell


"""Edit table is used to merge cells in table and add text to created table. Takes table object,
text and cell position as parameters"""


def edit_table(table, row, column, text=None, merge=False, to_row=0, to_column=0):
    if merge:
        table.cell(row, column).merge(table.cell(to_row, to_column))
    table.cell(row, column).text_frame.auto_size = MSO_AUTO_SIZE.SHAPE_TO_FIT_TEXT
    table.cell(row, column).text = text


"""Edit table style helps to format text present inside a table."""


def edit_table_style(table, size, font='Arial', text_alignment='center', bold=False, italic=False, underline=False,
                     color=(0, 0, 0)):
    text_alignment = text_alignment.lower()
    text_aligner = {'left': PP_PARAGRAPH_ALIGNMENT.LEFT, 'right': PP_PARAGRAPH_ALIGNMENT.RIGHT,
                    'center': PP_PARAGRAPH_ALIGNMENT.CENTER}
    for cell in iter_cells(table):
        for paragraph in cell.text_frame.paragraphs:
            for run in paragraph.runs:
                run.font.size = Pt(size)
                run.font.name = font.capitalize()
                run.font.bold = bold
                run.font.italic = italic
                run.font.underline = underline
                run.font.color.rgb = RGBColor(color[0], color[1], color[2])
            paragraph.alignment = text_aligner[text_alignment]

 
"""Helps yo add picture to ppt"""

def add_picture(slide,path,left=0,top=0,width=None,height=None):
    slide.shapes.add_picture(path,left,top,width=Inches(width),height=Inches(height))

"""Saves ppt at given path"""


def save_ppt(ppt, path):
    ppt.save(path)
